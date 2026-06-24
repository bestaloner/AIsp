#!/usr/bin/env python3
"""Generate VisionGuard A4 double-sided PPTX from HTML flyer content."""

from pptx import Presentation
from pptx.util import Cm, Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
import copy

# ── Constants ──────────────────────────────────────────
BLUE = RGBColor(0x25, 0x63, 0xEB)
BLUE_DARK = RGBColor(0x1D, 0x4E, 0xD8)
BLUE_LIGHT = RGBColor(0x3B, 0x82, 0xF6)
ORANGE = RGBColor(0xF5, 0x7A, 0x00)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
GRAY = RGBColor(0x64, 0x74, 0x8B)
GRAY_LIGHT = RGBColor(0x94, 0xA3, 0xB8)
TEXT = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
BG_LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
BG_BLUE = RGBColor(0xEF, 0xF6, 0xFF)

# 12 hazard colors
HZ_COLORS = [
    RGBColor(0xF8, 0x71, 0x71), RGBColor(0xFB, 0x92, 0x3C), RGBColor(0xFA, 0xCC, 0x15),
    RGBColor(0x4A, 0xDE, 0x80), RGBColor(0x60, 0xA5, 0xFA), RGBColor(0xA7, 0x8B, 0xFA),
    RGBColor(0xFB, 0x92, 0x3C), RGBColor(0xFB, 0xBF, 0x24), RGBColor(0x34, 0xD3, 0x99),
    RGBColor(0x81, 0x8C, 0xF8), RGBColor(0x9C, 0xA3, 0xAF), RGBColor(0x6B, 0x72, 0x80),
]
HZ_NAMES = [
    ("No Hard Hat", "安全帽未佩戴", "行为"), ("No Safety Vest", "反光衣未穿", "行为"),
    ("Edge Protection", "临边防护缺失", "物"), ("Smoking", "吸烟检测", "行为"),
    ("Intrusion", "人员越界", "行为"), ("Pit Overload", "坑边堆载", "物"),
    ("Open Flame", "明火检测", "行为"), ("Bulk Lifting", "散料吊装", "物"),
    ("Crowd Gather", "人员聚集", "行为"), ("Vehicle Violation", "车辆违规", "行为"),
    ("Suspended Platform", "吊篮作业", "物"), ("Enclosure Breach", "工地围挡封闭", "物"),
]
WF_STEPS = [
    ("AI Detection", "AI智能检测", RGBColor(0x25, 0x63, 0xEB), "📸"),
    ("Human Review", "人工复核", ORANGE, "👁"),
    ("WO Dispatch", "工单派发", RGBColor(0xE1, 0x1D, 0x48), "📋"),
    ("Rectification", "现场整改", GREEN, "🔧"),
    ("Verification", "整改验收", RGBColor(0x8B, 0x5C, 0xF6), "✅"),
    ("Closure", "闭环归档", RGBColor(0x37, 0x30, 0xA3), "🔒"),
]
SLIDE_W = Cm(21)
SLIDE_H = Cm(29.7)

# ── Helpers ────────────────────────────────────────────
def add_textbox(slide, left, top, width, height, text="", font_size=Pt(8),
                color=TEXT, bold=False, align=PP_ALIGN.LEFT, font_name='Segoe UI',
                line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
        spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
        lnSpc.append(spcPct)
        pPr.append(lnSpc)
    return txBox

def add_multiline_box(slide, left, top, width, height, lines, font_name='Segoe UI'):
    """lines: list of (text, font_size, color, bold, align)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        text, fsize, color, bold, align = ln
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = fsize
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(0)
        p.space_before = Pt(1)
    return txBox

def add_rect(slide, left, top, width, height, fill=None, border_color=BORDER,
             border_width=Pt(1), corner_radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
                                   left, top, width, height)
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape

def add_img_placeholder(slide, left, top, width, height, label, desc=""):
    """Dashed-border placeholder for images."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    shape.line.color.rgb = GRAY_LIGHT
    shape.line.width = Pt(1)
    # Set dashed line
    shape.line.dash_style = 2  # dash
    # Add text
    tf = shape.text_frame
    tf.word_wrap = True
    # Center vertically
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Icon
    p = tf.paragraphs[0]
    p.text = "🖼"
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(2)
    # Label
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(7)
    p2.font.color.rgb = GRAY
    p2.font.bold = True
    p2.font.name = 'Segoe UI'
    p2.alignment = PP_ALIGN.CENTER
    if desc:
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(6)
        p3.font.color.rgb = GRAY_LIGHT
        p3.font.name = 'Segoe UI'
        p3.alignment = PP_ALIGN.CENTER
    return shape

def add_slide_header(slide, title="VisionGuard", subtitle="AI Construction Safety Inspection Platform · 建筑施工AI安全巡检平台",
                     badge="AI-POWERED · YOLO + VLM"):
    """Add standard header bar to a slide."""
    # Blue line at top
    add_rect(slide, Cm(0), Cm(0), SLIDE_W, Cm(0.15), fill=BLUE, border_color=BLUE)
    # Logo square
    add_rect(slide, Cm(1), Cm(0.5), Cm(1.3), Cm(1.3), fill=BLUE, corner_radius=Cm(0.25))
    add_textbox(slide, Cm(1), Cm(0.5), Cm(1.3), Cm(1.3), "🛡", Pt(12), WHITE, True, PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, Cm(2.6), Cm(0.5), Cm(8), Cm(0.7), title, Pt(16), BLUE, True)
    add_textbox(slide, Cm(2.6), Cm(1.15), Cm(10), Cm(0.4), subtitle, Pt(6), GRAY)
    # Badge
    badge_box = add_rect(slide, Cm(17), Cm(0.65), Cm(3.2), Cm(0.7), fill=BLUE, corner_radius=Cm(0.3))
    add_textbox(slide, Cm(17), Cm(0.65), Cm(3.2), Cm(0.7), badge, Pt(5.5), WHITE, True, PP_ALIGN.CENTER)
    # Separator line
    add_rect(slide, Cm(1), Cm(2.05), Cm(19), Cm(0.03), fill=BORDER, border_color=BORDER)

def add_slide_footer(slide):
    """Add standard footer."""
    add_rect(slide, Cm(1), Cm(28.4), Cm(19), Cm(0.03), fill=BORDER, border_color=BORDER)
    add_textbox(slide, Cm(1), Cm(28.6), Cm(8), Cm(0.5),
                "VisionGuard · AI Construction Safety · © 2026", Pt(5.5), GRAY)
    add_textbox(slide, Cm(11), Cm(28.6), Cm(9), Cm(0.5),
                "🌐 Live Demo: bestaloner.github.io/AIsp  |  📧 contact@visionguard.ai",
                Pt(5.5), GRAY, align=PP_ALIGN.RIGHT)

def add_section_title(slide, left, top, width, text, dot_color=BLUE):
    """Add a section title with colored dot and underline."""
    # Dot
    add_rect(slide, left, top, Cm(0.45), Cm(0.45), fill=dot_color, border_color=dot_color, corner_radius=Cm(0.1))
    # Title
    add_textbox(slide, left + Cm(0.7), top - Cm(0.02), width, Cm(0.5), text, Pt(8), TEXT, True)
    # Underline
    add_rect(slide, left, top + Cm(0.55), width, Cm(0.02), fill=BORDER, border_color=BORDER)

# ═══════════════════════════════════════════════════════
#  CREATE PRESENTATION
# ═══════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]  # blank

# ───────────────────────────────────────────────────────
#  SLIDE 1: FRONT PAGE — Product Overview
# ───────────────────────────────────────────────────────
slide1 = prs.slides.add_slide(blank_layout)
# White background is default

add_slide_header(slide1)

# ── Hero section ──
y_hero = Cm(2.3)
# Hero text
add_textbox(slide1, Cm(1), y_hero, Cm(11), Cm(1.6),
            "Your Site's AI Safety Net That Never Sleeps",
            Pt(18), BLUE_DARK, True)
add_textbox(slide1, Cm(1), y_hero + Cm(1.5), Cm(11), Cm(1.2),
            "VisionGuard combines real-time YOLO object detection with VLM secondary verification to monitor "
            "construction sites 24/7. AI flags every potential hazard — human reviewers confirm or dismiss — "
            "and the platform tracks every issue through to rectification closure.",
            Pt(7), GRAY)
add_textbox(slide1, Cm(1), y_hero + Cm(2.5), Cm(11), Cm(0.5),
            "全天候AI视频智能巡检 · YOLO实时检测+VLM大模型二次复核 · 人机协同精准判识 · 全链路闭环管理",
            Pt(6.5), BLUE, True)

# Metrics row
metrics = [("12", "Hazard Types\n隐患类型"), ("4", "User Roles\n角色权限(RBAC)"),
           ("6", "Closed-Loop\n全流程闭环"), ("24/7", "Real-Time\n全天候监控")]
for i, (num, label) in enumerate(metrics):
    mx = Cm(1 + i * 2.8)
    box = add_rect(slide1, mx, y_hero + Cm(3.1), Cm(2.5), Cm(1.3), fill=BG_BLUE, corner_radius=Cm(0.2))
    add_textbox(slide1, mx, y_hero + Cm(3.15), Cm(2.5), Cm(0.7), num, Pt(14), BLUE, True, PP_ALIGN.CENTER)
    add_textbox(slide1, mx, y_hero + Cm(3.85), Cm(2.5), Cm(0.6), label, Pt(5.5), GRAY, align=PP_ALIGN.CENTER)

# Image placeholders (right side)
add_img_placeholder(slide1, Cm(12.5), y_hero, Cm(7.5), Cm(2.1),
                    "AI Review Workbench", "人工复核工作台截图")
add_img_placeholder(slide1, Cm(12.5), y_hero + Cm(2.3), Cm(7.5), Cm(2.1),
                    "Executive Dashboard", "工作台总览截图")

# ── 12 Hazard Types ──
y_hz = Cm(7.2)
add_section_title(slide1, Cm(1), y_hz, Cm(8), "12 AI Detection Capabilities / AI检测能力覆盖", BLUE)

# 6×2 grid
for i, (en, cn, tag) in enumerate(HZ_NAMES):
    col = i % 6
    row = i // 6
    hx = Cm(1 + col * 3.15)
    hy = y_hz + Cm(0.6) + row * Cm(0.8)
    # Color dot
    add_rect(slide1, hx, hy + Cm(0.1), Cm(0.35), Cm(0.35), fill=HZ_COLORS[i],
             border_color=HZ_COLORS[i], corner_radius=Cm(0.08))
    # Name
    add_textbox(slide1, hx + Cm(0.5), hy, Cm(2.2), Cm(0.35), en, Pt(6.5), TEXT, True)
    add_textbox(slide1, hx + Cm(0.5), hy + Cm(0.32), Cm(1.2), Cm(0.25), cn, Pt(5), GRAY_LIGHT)
    # Tag
    tag_color = RGBColor(0x1E, 0x40, 0xAF) if tag == "行为" else RGBColor(0x92, 0x40, 0x0E)
    tag_bg = RGBColor(0xDB, 0xEA, 0xFE) if tag == "行为" else RGBColor(0xFE, 0xF3, 0xC7)
    tag_box = add_rect(slide1, hx + Cm(2), hy + Cm(0.05), Cm(0.7), Cm(0.45),
                       fill=tag_bg, corner_radius=Cm(0.3))
    add_textbox(slide1, hx + Cm(2), hy + Cm(0.05), Cm(0.7), Cm(0.45), tag, Pt(4.5), tag_color, align=PP_ALIGN.CENTER)

# ── 6-Step Workflow ──
y_wf = Cm(9.3)
add_section_title(slide1, Cm(1), y_wf, Cm(8), "6-Step Closed-Loop Workflow / 六步全流程闭环", ORANGE)

for i, (en, cn, color, icon) in enumerate(WF_STEPS):
    wx = Cm(1 + i * 3.15)
    # Circle with icon
    circle = add_rect(slide1, wx + Cm(0.5), y_wf + Cm(0.6), Cm(1.1), Cm(1.1),
                      fill=color, border_color=color, corner_radius=Cm(0.55))
    add_textbox(slide1, wx + Cm(0.5), y_wf + Cm(0.65), Cm(1.1), Cm(1.0), icon, Pt(12), WHITE, align=PP_ALIGN.CENTER)
    # Labels
    add_textbox(slide1, wx, y_wf + Cm(1.75), Cm(2.15), Cm(0.35), en, Pt(6.5), TEXT, True, PP_ALIGN.CENTER)
    add_textbox(slide1, wx, y_wf + Cm(2.05), Cm(2.15), Cm(0.25), cn, Pt(5), GRAY, align=PP_ALIGN.CENTER)
    # Arrow between circles (except last)
    if i < 5:
        add_textbox(slide1, wx + Cm(1.85), y_wf + Cm(0.85), Cm(1), Cm(0.5),
                    "▸", Pt(10), GRAY_LIGHT, align=PP_ALIGN.CENTER)

# ── Bottom image placeholders ──
y_bot = Cm(11.8)
add_img_placeholder(slide1, Cm(1), y_bot, Cm(6.2), Cm(2.8),
                    "Camera Live View + AI Annotations", "摄像头实时画面 + AI标注效果示意")
add_img_placeholder(slide1, Cm(7.5), y_bot, Cm(6.2), Cm(2.8),
                    "Hazard Detection Examples", "隐患检测示例：安全帽未佩戴、临边防护缺失")
add_img_placeholder(slide1, Cm(14), y_bot, Cm(6), Cm(2.8),
                    "Work Order Lifecycle Tracking", "工单全生命周期状态流转示意")

add_slide_footer(slide1)

# ───────────────────────────────────────────────────────
#  SLIDE 2: BACK PAGE — Platform Details
# ───────────────────────────────────────────────────────
slide2 = prs.slides.add_slide(blank_layout)

add_slide_header(slide2, badge="HUMAN-AI COLLABORATION")

y2 = Cm(2.3)

# ── Section: Core Platform Modules ──
add_section_title(slide2, Cm(1), y2, Cm(10), "Core Platform Modules / 核心功能模块", GREEN)

# Feature blocks (3 rows, each with thumbnail + text)
features = [
    ("Executive Dashboard & AI Review Workbench",
     "工作台总览 · 人工复核工作台",
     "Real-time KPI cards, hazard trend charts, and project status overview. The AI Review Workbench "
     "presents original camera images side-by-side with AI-annotated detection results — reviewers "
     "confirm or dismiss each finding in seconds. This human-in-the-loop step combines AI speed "
     "with expert judgment to eliminate false positives.",
     "Dashboard + AI Review", "工作台 + 复核台截图"),
    ("Work Order Dispatch & Lifecycle Tracking",
     "工单派发 · 工单台账",
     "One-click batch dispatch confirmed hazards as work orders to assigned on-site executors. "
     "Full lifecycle: Pending → Dispatched → Rectifying → Submitted → Verified → Closed. "
     "Every status change is timestamped and logged. Executors see only assigned tasks; "
     "admins have full visibility across all projects.",
     "Work Order Management", "工单派发 + 工单台账截图"),
    ("Multi-Site Management & Per-Camera Scene Config",
     "项目管理 · 摄像头场景矩阵配置 · 工单设置",
     "One platform manages unlimited construction sites. Each camera independently configured "
     "with a scene matrix — enable only relevant hazard types per camera's field of view. "
     "WO Settings control which hazard categories auto-generate work orders. "
     "Project-level executor assignment ensures accountability.",
     "Camera & Project Config", "摄像头配置 + 项目管理截图"),
]

for fi, (title, cn, desc, thumb_label, thumb_desc) in enumerate(features):
    fy = y2 + Cm(0.6) + fi * Cm(3.05)
    # Thumbnail placeholder
    add_img_placeholder(slide2, Cm(1), fy, Cm(6), Cm(2.7), thumb_label, thumb_desc)
    # Text
    add_textbox(slide2, Cm(7.4), fy, Cm(12.6), Cm(0.5), title, Pt(8), TEXT, True)
    add_textbox(slide2, Cm(7.4), fy + Cm(0.55), Cm(12.6), Cm(0.4), cn, Pt(6), GRAY)
    add_textbox(slide2, Cm(7.4), fy + Cm(1.0), Cm(12.6), Cm(1.6), desc, Pt(6.5), GRAY)

# ── Roles + Advantages (below features) ──
y3 = Cm(11.9)

# Left: Roles
add_section_title(slide2, Cm(1), y3, Cm(8), "Role-Based Access / 角色权限", BLUE)

roles = [
    ("👑", "Administrator", "系统管理员", "Full platform access\nUser/Project/Camera mgmt", BG_BLUE, BLUE_DARK),
    ("🔍", "Reviewer", "复核员", "AI review workbench\nConfirm/dismiss hazards", RGBColor(0xFE,0xF3,0xC7), ORANGE),
    ("🔧", "Executor", "执行人", "Receives assigned WOs\nOn-site rectification", RGBColor(0xD1,0xFA,0xE5), GREEN),
    ("👀", "Viewer", "观察员", "Read-only access\nDashboard & reports", RGBColor(0xF8,0xFA,0xFC), GRAY),
]
for i, (icon, en, cn, caps, bg, accent) in enumerate(roles):
    rx = Cm(1 + i * 2.45)
    ry = y3 + Cm(0.6)
    add_rect(slide2, rx, ry, Cm(2.15), Cm(3.0), fill=WHITE, border_color=BORDER, corner_radius=Cm(0.2))
    # Icon circle
    add_rect(slide2, rx + Cm(0.55), ry + Cm(0.2), Cm(1.0), Cm(1.0), fill=bg, border_color=bg, corner_radius=Cm(0.5))
    add_textbox(slide2, rx + Cm(0.55), ry + Cm(0.25), Cm(1.0), Cm(0.9), icon, Pt(10), accent, align=PP_ALIGN.CENTER)
    # Text
    add_textbox(slide2, rx + Cm(0.1), ry + Cm(1.4), Cm(1.95), Cm(0.35), en, Pt(6.5), TEXT, True, PP_ALIGN.CENTER)
    add_textbox(slide2, rx + Cm(0.1), ry + Cm(1.7), Cm(1.95), Cm(0.25), cn, Pt(5), GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide2, rx + Cm(0.1), ry + Cm(2.05), Cm(1.95), Cm(0.8), caps, Pt(5.5), GRAY, align=PP_ALIGN.CENTER)

# Right: Advantages
adv_x = Cm(11.5)
add_section_title(slide2, adv_x, y3, Cm(8), "Why VisionGuard / 核心优势", ORANGE)

advantages = [
    ("Dual AI Engine / 双AI引擎", "YOLO real-time detection + VLM secondary verification — high recall, low false positives"),
    ("Human-in-the-Loop / 人机协同", "AI flags, humans decide — combining machine speed with expert judgment for unmatched accuracy"),
    ("Full Lifecycle Traceability / 全链路可追溯", "Every hazard has a complete audit trail from AI detection to rectification closure — compliance ready"),
    ("Enterprise Multi-Site Scale / 企业级多项目", "One platform manages all sites, cameras, and teams — scales from single site to enterprise portfolio"),
]
for ai, (title, desc) in enumerate(advantages):
    ay = y3 + Cm(0.6) + ai * Cm(0.7)
    # Check mark
    add_rect(slide2, adv_x, ay, Cm(0.4), Cm(0.4), fill=RGBColor(0xD1, 0xFA, 0xE5),
             border_color=RGBColor(0xD1, 0xFA, 0xE5), corner_radius=Cm(0.2))
    add_textbox(slide2, adv_x, ay - Cm(0.02), Cm(0.4), Cm(0.4), "✓", Pt(5), GREEN, align=PP_ALIGN.CENTER)
    # Content
    add_textbox(slide2, adv_x + Cm(0.65), ay - Cm(0.05), Cm(7.5), Cm(0.3), title, Pt(6.5), TEXT, True)
    add_textbox(slide2, adv_x + Cm(0.65), ay + Cm(0.25), Cm(7.5), Cm(0.4), desc, Pt(5.5), GRAY)

# ── Tech + Use Case ──
y4 = Cm(15.1)

# Left: Tech Stack
add_section_title(slide2, Cm(1), y4, Cm(8), "Technology Stack / 技术栈", GRAY)

techs = [("🎯 YOLO", "Real-time Object Detection"), ("🧠 VLM", "Vision Language Model (Qwen)"),
         ("⚡ vLLM", "High-Throughput Inference"), ("🔗 REST API", "Upstream Integration")]
for i, (name, sub) in enumerate(techs):
    tx = Cm(1 + i * 2.45)
    add_rect(slide2, tx, y4 + Cm(0.65), Cm(2.15), Cm(1.0), fill=BG_LIGHT, corner_radius=Cm(0.15))
    add_textbox(slide2, tx, y4 + Cm(0.7), Cm(2.15), Cm(0.45), name, Pt(7), TEXT, True, PP_ALIGN.CENTER)
    add_textbox(slide2, tx, y4 + Cm(1.1), Cm(2.15), Cm(0.3), sub, Pt(5), GRAY, align=PP_ALIGN.CENTER)

# Data pipeline
add_textbox(slide2, Cm(1), y4 + Cm(1.85), Cm(9.5), Cm(0.5),
            "📸 Cameras → 🎯 YOLO → 🧠 VLM → 👁 Human Review → 📋 Work Order → 🔧 Rectify → ✅ Verify → 📤 Upstream API",
            Pt(6), BLUE_DARK, True)
add_textbox(slide2, Cm(1), y4 + Cm(2.15), Cm(9.5), Cm(0.4),
            "指数退避重试 + 死信队列确保数据零丢失  ·  Exponential backoff retry + dead-letter queue for zero data loss",
            Pt(5.5), GRAY)

# Right: Use Case
uc_x = Cm(11.5)
add_section_title(slide2, uc_x, y4, Cm(8), "Application Scenario / 应用场景", GREEN)

add_rect(slide2, uc_x, y4 + Cm(0.65), Cm(8.5), Cm(2.1), fill=BG_BLUE, corner_radius=Cm(0.25))
add_img_placeholder(slide2, uc_x + Cm(0.3), y4 + Cm(0.8), Cm(3.5), Cm(1.8),
                    "Site Overview", "工地全景示意")
add_textbox(slide2, uc_x + Cm(4.1), y4 + Cm(0.8), Cm(4.1), Cm(0.4),
            "Typical Deployment / 典型部署", Pt(7), BLUE_DARK, True)
add_textbox(slide2, uc_x + Cm(4.1), y4 + Cm(1.2), Cm(4.1), Cm(1.5),
            "A mid-size construction site deploys 8–20 cameras covering entrances, edges, lifting zones, "
            "material yards, and access roads. VisionGuard monitors all feeds simultaneously — "
            "impossible for human guards to match. When a hazard is detected, the system alerts "
            "reviewers within seconds and dispatches work orders to on-site executors.\n\n"
            "中型工地通常部署8-20路摄像头覆盖出入口、临边、吊装区、材料堆场等关键区域。AI秒级告警、"
            "人工快速复核、工单直达现场执行人。",
            Pt(5.5), GRAY)

add_slide_footer(slide2)

# ── Save ──────────────────────────────────────────────
output_path = r"f:\0.AI设计库\ai视频识别\gh-pages-deploy\VisionGuard_Product_Overview.pptx"
prs.save(output_path)
print(f"PPTX saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Slide size: {prs.slide_width/914400:.1f}cm x {prs.slide_height/914400:.1f}cm")
